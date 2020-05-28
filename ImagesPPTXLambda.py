from UtilsLambda import FindNewFilename, NewFilePath
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

#1px =~ 9525EMU

#Ensure .jpgs & rotate landscapes
def ProcessImages(TEMPDIR, image_ext):
    for file in next(os.walk(TEMPDIR))[2]: #files in TEMPDIR
        FName, FExt = os.path.splitext(file)
        if FExt in image_ext: #if file is image
            if not FExt.lower() == '.jpg': #if file is not jpg
                ConvertToJpg(TEMPDIR, file)

            #Check orientation
            RotateToPortrait(os.path.join(TEMPDIR, file))
            #allPageDimensions.update(GetImageDimensionsInches(os.path.join(TEMPDIR, file))) #should add new k&v to dict
            #width, height = GetImageDimensionsInches((os.path.join(TEMPDIR, file)))
            # if width > maxwidth:
            #     maxwidth = width
            # if height > maxheight:
            #     maxheight = height
    #print("maxheight " + str(maxheight) + "maxwidth " + str(maxwidth))
    return True

#Convert images to jpg
def ConvertToJpg(Destination, File): # remember to pass TEMPDIR + FILE
    comicPage = Image.open(File)
    newComicPage = FindNewFilename(Destination, File) #to make sure a page is not overwritten
    comicPage.convert('RGB').save(newComicPage, quality=95) #save
    os.unlink(os.path.join(File)) #delete original

#Rotate Images to portrait
def RotateToPortrait(File):
    comicPage = Image.open(File)
    #Check page is portrait
    width, height = comicPage.size

    if width > height:
        #print(width, height)
        #print('Rotating file: ' + File)
        # Angle is in degrees counter clockwise
        comicPage = comicPage.rotate(270, expand=True)
        # crop the rotated image to the size of the original image
        comicPage.save(File, dpi=(300, 300), quality=95)
        width, height = comicPage.size
        #print(width, height)
        comicPage.close()

#Returns dimensions for first .jpg in folder
def FirstImageDimensions(TEMPDIR):
    for file in next(os.walk(TEMPDIR))[2]:
        FName, FExt = os.path.splitext(file)
        if FExt.lower() == '.jpg':
            width, height = GetImageDimensionsInches(os.path.join(TEMPDIR, file))
            return width, height

#Returns dimensions for provided file in inches
def GetImageDimensionsInches(File):
    comicPage = Image.open(File)
    width, height = 0, 0
    try:
        if comicPage.info['dpi'] != None:
            width, height = comicPage.size
            width = width / comicPage.info['dpi'][0] #gives inches
            height = height / comicPage.info['dpi'][0]  # gives inches
            #print('Image DPI: ' + str(comicPage.info['dpi'][0]))
    except KeyError:
        #print('DPI info not found, assuming 300')
        width, height = comicPage.size
        width = width / 300  # gives inches
        height = height / 300 # gives inches
    return width, height

def MakePresentation(width, height):
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return prs

def AddSlide(prs, file):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    top = 0
    left = 0
    pic = slide.shapes.add_picture(file, left, top, width = prs.slide_width, height = prs.slide_height)
    return prs

def AddXmlSlide(prs, XmlDict):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    cols = 2
    rows = len(XmlDict)
    left = top = int(prs.slide_width * .05)
    width = int(prs.slide_width * .9)
    height = 1

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = int(width * .2)
    table.columns[1].width = int(width * .8)
    row = 0 #start at one because title
    for key in XmlDict:
        if isinstance(XmlDict[key], str): #if value is a string

            print(key, XmlDict[key])
            table.cell(row, 0).text = key
            table.cell(row, 1).text = str(XmlDict[key])
            row += 1

    for r in range(len(XmlDict)):
        for c in range(0, 2):
            cell = table.cell(r, c)
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(6)
            #para.font.name = 'Comic Sans MS'
            para.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
    return prs

def SavePPTX(prs, filename, TEMPDIR):
    newFile = NewFilePath(os.path.join(TEMPDIR, filename))
    prs.save(newFile)
    print("New Comic Created: " + newFile)
    return newFile

