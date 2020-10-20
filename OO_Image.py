from PIL import Image

import os
import time
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor, ColorFormat
import warnings

# from pptx.enum.dml import MSO_COLOR_TYPE # MSO_THEME_COLOR
# import codecs

# 1px =~ 9525EMU

class CS_Image:
    def __init__(self, SUBMITTED_FILE="", TEMPDIR=""):
        self.SUBMITTED_FILE = SUBMITTED_FILE # full path and file
        self.TEMPDIR = TEMPDIR

        # Some image EXIF data causes errors. This hides that warning.
        warnings.filterwarnings("ignore", "(Possibly )?corrupt EXIF data", UserWarning)

        """ END OF INIT """

    # Ensure .jpgs & rotate landscapes
    def process_images(TEMPDIR, image_ext):
        for file in next(os.walk(TEMPDIR))[2]:  # files in TEMPDIR
            FName, FExt = os.path.splitext(file)
            if FExt in image_ext:  # if file is image
                if not FExt.lower() == '.jpg':  # if file is not jpg
                    convert_to_jpg(TEMPDIR, file)

                # Check orientation
                rotate_to_portrait(os.path.join(TEMPDIR, file))
                # allPageDimensions.update(GetImageDimensionsInches(os.path.join(TEMPDIR, file))) #should add new k&v to dict
                # width, height = GetImageDimensionsInches((os.path.join(TEMPDIR, file)))
                # if width > maxwidth:
                #     maxwidth = width
                # if height > maxheight:
                #     maxheight = height
        # print("maxheight " + str(maxheight) + "maxwidth " + str(maxwidth))
        return True

    # Convert images to jpg
    def convert_to_jpg(self, file, output_file=""):  # remember to pass TEMPDIR + FILE
        if output_file == "":
            output_file == file
            print("File is " + file)
        comicPage = Image.open(file)
        print("outputFile is " + output_file)
        comicPage.convert('RGB').save(output_file, quality=95)  # save
        print("File conversion complete")
        # os.unlink(os.path.join(self.TEMPDIR, file))  # delete original

    # Rotate Images to portrait
    def rotate_to_portrait(self, file):
        comicPage = Image.open(file)
        # Check page is portrait
        width, height = comicPage.size

        if width > height:
            print("Rotating image: " + file)
            # print(width, height)
            # print('Rotating file: ' + File)
            # Angle is in degrees counter clockwise
            comicPage = comicPage.rotate(270, expand=True)
            # crop the rotated image to the size of the original image
            comicPage.save(file, dpi=(300, 300), quality=95)
            width, height = comicPage.size
            # print(width, height)
            comicPage.close()

    # Returns dimensions for first .jpg in folder
    def first_image_dimensions(self, TEMPDIR):
        for file in next(os.walk(TEMPDIR))[2]:
            FName, FExt = os.path.splitext(file)
            if FExt.lower() == '.jpg':
                width, height = GetImageDimensionsInches(os.path.join(TEMPDIR, file))
                print(file, width, height)
                return width, height

    # Returns dimensions for provided file in inches
    def get_image_dimensions_inches(self, file):
        comicPage = Image.open(file)
        width, height = 0, 0
        try:
            if comicPage.info['dpi'] != None:
                width, height = comicPage.size
                width = width / comicPage.info['dpi'][0]  # gives inches
                height = height / comicPage.info['dpi'][0]  # gives inches
                # print('Image DPI: ' + str(comicPage.info['dpi'][0]))
        except KeyError:
            # print('DPI info not found, assuming 300')
            width, height = comicPage.size
            width = width / 300  # gives inches
            height = height / 300  # gives inches
        return width, height

    def make_presentation(self, width, height):
        prs = Presentation()
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)
        return prs

    def add_slide(self, prs, file):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        top = 0
        left = 0
        pic = slide.shapes.add_picture(file, left, top, width=prs.slide_width, height=prs.slide_height)
        return prs

    def add_xml_slide(self, prs, xml_dict):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0, 0, 0)

        cols = 2
        rows = len(xml_dict)
        left = top = int(prs.slide_width * .05)
        width = int(prs.slide_width * .9)
        height = 1

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = int(width * .2)
        table.columns[1].width = int(width * .8)
        row = 0  # start at one because title
        for key in xml_dict:
            if isinstance(xml_dict[key], str):  # if value is a string

                print(key, xml_dict[key])
                table.cell(row, 0).text = key
                table.cell(row, 1).text = str(xml_dict[key])
                row += 1

        for r in range(len(xml_dict)):
            for c in range(0, 2):
                cell = table.cell(r, c)
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(6)
                # para.font.name = 'Comic Sans MS'
                para.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
        return prs

    def save_pptx(self, prs, new_pptx_folder_and_filename):

        if os.path.isfile(new_pptx_folder_and_filename) == True:
            print('Old file found. Deleting: ' + new_pptx_folder_and_filename)
            os.unlink(new_pptx_folder_and_filename)
            time.sleep(2)
        prs.save(new_pptx_folder_and_filename)
        print("New Comic Created: " + new_pptx_folder_and_filename)